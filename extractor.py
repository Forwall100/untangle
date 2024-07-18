import os
from typing import Callable, Dict
import zipfile
import xml.etree.ElementTree as ET
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from PyPDF2 import PdfReader
import csv
import chardet
import json
import yaml
import markdown
from bs4 import BeautifulSoup
import email
import ebooklib
from ebooklib import epub
from mobi import extract as mobi_extract
import html2text
import tempfile
import requests
from urllib.parse import urlparse


def extract_epub(file_path: str) -> str:
    book = epub.read_epub(file_path)
    text = []
    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            text.append(html2text.html2text(item.get_content().decode("utf-8")))
    return "\n".join(text)


def extract_fb2(file_path: str) -> str:
    with open(file_path, "rb") as file:
        content = file.read()
    soup = BeautifulSoup(content, "xml")
    return " ".join(soup.stripped_strings)


def extract_docx(file_path: str) -> str:
    doc = Document(file_path)
    return "\n".join(paragraph.text for paragraph in doc.paragraphs)


def extract_xlsx(file_path: str) -> str:
    wb = load_workbook(file_path)
    return "\n".join(
        str(cell.value)
        for sheet in wb.worksheets
        for row in sheet.rows
        for cell in row
        if cell.value
    )


def extract_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    return "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )


def extract_pdf(file_path: str) -> str:
    reader = PdfReader(file_path)
    return "\n".join(page.extract_text() for page in reader.pages)


def extract_opendocument(file_path: str) -> str:
    with zipfile.ZipFile(file_path) as zf:
        content = zf.read("content.xml")
    root = ET.fromstring(content)
    return "\n".join(elem.text for elem in root.iter() if elem.text)


def extract_text_with_encoding(file_path: str) -> str | None:
    with open(file_path, "rb") as file:
        raw_data = file.read()

    chardet_data = chardet.detect(raw_data)
    encoding = chardet_data.get("encoding")

    if encoding is None:
        return None

    return raw_data.decode(encoding)


def extract_csv(file_path: str) -> str:
    with open(file_path, "r", newline="", encoding="utf-8") as csvfile:
        reader = csv.reader(csvfile)
        return "\n".join(",".join(row) for row in reader)


def extract_json(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as file:
        return json.dumps(json.load(file), indent=2)


def extract_yaml(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as file:
        return yaml.dump(yaml.safe_load(file), default_flow_style=False)


def extract_markdown(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as file:
        return markdown.markdown(file.read())


def extract_html(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as file:
        soup = BeautifulSoup(file, "html.parser")
        return soup.get_text()


def extract_xml(file_path: str) -> str:
    tree = ET.parse(file_path)
    return ET.tostring(tree.getroot(), encoding="unicode", method="text")


def extract_eml(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as file:
        msg = email.message_from_file(file)
        return f"Subject: {msg['subject']}\nFrom: {msg['from']}\nTo: {msg['to']}\n\n{msg.get_payload()}"


def is_url(path: str) -> bool:
    try:
        result = urlparse(path)
        return all([result.scheme, result.netloc])
    except ValueError:
        return False


def extract_webpage(url: str) -> str:
    try:
        response = requests.get(url)
        response.raise_for_status()
        soup = BeautifulSoup(response.content, "html.parser")

        for script in soup(["script", "style"]):
            script.decompose()

        text = soup.get_text()

        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = "\n".join(chunk for chunk in chunks if chunk)

        return text
    except requests.RequestException as e:
        return f"Ошибка при получении веб-страницы: {str(e)}"


def extract_unknown(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
            content = file.read()
        if content:
            return content.strip()
        else:
            text = extract_text_with_encoding(file_path)
            if text is not None:
                return text.strip()
            else:
                return "Не удалось извлечь текст из файла"
    except Exception as e:
        return f"Не удалось прочитать файл: {str(e)}"


EXTRACTORS = {
    ".docx": extract_docx,
    ".doc": extract_docx,
    ".xlsx": extract_xlsx,
    ".xls": extract_xlsx,
    ".pptx": extract_pptx,
    ".ppt": extract_pptx,
    ".pdf": extract_pdf,
    ".odt": extract_opendocument,
    ".ods": extract_opendocument,
    ".odp": extract_opendocument,
    ".txt": extract_text_with_encoding,
    ".rtf": extract_text_with_encoding,
    ".csv": extract_csv,
    ".json": extract_json,
    ".yml": extract_yaml,
    ".yaml": extract_yaml,
    ".md": extract_markdown,
    ".html": extract_html,
    ".htm": extract_html,
    ".xml": extract_xml,
    ".eml": extract_eml,
    ".log": extract_text_with_encoding,
    ".ini": extract_text_with_encoding,
    ".epub": extract_epub,
    ".fb2": extract_fb2,
}


def extract_text(path: str) -> str:
    if is_url(path):
        return extract_webpage(path)

    _, ext = os.path.splitext(path)
    ext = ext.lower()

    try:
        extractor = EXTRACTORS.get(ext, extract_unknown)
        return extractor(path)
    except Exception as e:
        return f"Ошибка при обработке файла: {str(e)}"
